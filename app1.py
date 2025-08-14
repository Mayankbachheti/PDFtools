import os
from flask import Flask, request, send_file, jsonify
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from pdf2docx import Converter 
from flask_cors import CORS
from PIL import Image
import zipfile

app = Flask(__name__)
CORS(app)

# Define the upload folder path
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'jpg', 'jpeg', 'png'}

# Ensure the upload folder exists
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Function to check allowed file extensions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def home():
    return "Welcome to the PDF Conversion and Compression API!"

# PDF to Word conversion route
@app.route('/convert_pdf_to_word', methods=['POST'])
def convert_pdf_to_word():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file and allowed_file(file.filename):
        # Save the file
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Convert PDF to Word
        output_word_path = file_path.rsplit('.', 1)[0] + '.docx'
        cv = Converter(file_path)
        cv.convert(output_word_path, start=0, end=None)  # Converts PDF to Word
        cv.close()

        return send_file(output_word_path, as_attachment=True)

    return jsonify({"error": "Invalid file format. Only PDF files are allowed."}), 400

# Compress PDF route
@app.route('/compress_pdf', methods=['POST'])
def compress_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Compress the PDF
        compressed_pdf_path = file_path.rsplit('.', 1)[0] + '_compressed.pdf'

        # Placeholder for PDF compression code
        # (Use a library like PyMuPDF or pikepdf to compress the PDF)

        # For now, we'll just mock the compression by copying the file
        with open(file_path, 'rb') as f_in:
            with open(compressed_pdf_path, 'wb') as f_out:
                f_out.write(f_in.read())

        return send_file(compressed_pdf_path, as_attachment=True)

    return jsonify({"error": "Invalid file format. Only PDF files are allowed."}), 400

# Compress image route (supports jpg, jpeg, png)
@app.route('/compress_image', methods=['POST'])
def compress_image():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Compress the image using Pillow
        img = Image.open(file_path)
        compressed_image_path = file_path.rsplit('.', 1)[0] + '_compressed.' + img.format.lower()

        # Resize or compress image (optional: adjust quality)
        img = img.convert("RGB")
        img.save(compressed_image_path, quality=85, optimize=True)

        return send_file(compressed_image_path, as_attachment=True)

    return jsonify({"error": "Invalid file format. Only image files (jpg, jpeg, png) are allowed."}), 400

# PDF to PowerPoint route (mock implementation)
'''@app.route('/convert_pdf_to_ppt', methods=['POST'])
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)

        # Placeholder for PDF to PowerPoint conversion
        # (You can use libraries like python-pptx and PyMuPDF to extract PDF content and create PPT slides)
        ppt_output_path = file_path.rsplit('.', 1)[0] + '.pptx'

        # For now, just mock the conversion by creating an empty PowerPoint file
        from pptx import Presentation
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # Add a blank slide
        slide.shapes.title.text = "Converted from PDF"

        ppt.save(ppt_output_path)

        return send_file(ppt_output_path, as_attachment=True)

    return jsonify({"error": "Invalid file format. Only PDF files are allowed."}), 400'''

if __name__ == '__main__':
    app.run(debug=True)
